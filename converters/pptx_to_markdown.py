"""
PPTX -> Document IR -> Markdown 변환기.

이미지 자체의 OCR/의미 분석은 범위 밖이며, 이미지는 메타데이터(도형 이름/크기)만
보존한다. Slide title, 텍스트, 표, 슬라이드 노트, 슬라이드 번호를 최대한 보존한다.
"""
from __future__ import annotations

import os
from typing import List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .document_ir import DocumentIR, ImageRef, SlideIR, TableIR


def pptx_to_ir(path: str) -> DocumentIR:
    prs = Presentation(path)
    slides: List[SlideIR] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = None
        paragraphs: List[str] = []
        tables: List[TableIR] = []
        images: List[ImageRef] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                is_title_placeholder = getattr(shape, "is_placeholder", False) and getattr(
                    shape.placeholder_format, "idx", None
                ) == 0
                if title is None and is_title_placeholder:
                    title = text
                else:
                    paragraphs.append(text)
            elif shape.has_table:
                table = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                headers = rows[0] if rows else []
                data_rows = rows[1:] if len(rows) > 1 else []
                tables.append(TableIR(headers=headers, rows=data_rows))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append(
                    ImageRef(shape_name=shape.name, width_emu=shape.width, height_emu=shape.height)
                )

        notes = None
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            notes = notes_text or None

        slides.append(
            SlideIR(
                slide_number=idx,
                title=title,
                paragraphs=paragraphs,
                tables=tables,
                notes=notes,
                images=images,
            )
        )

    doc_title = slides[0].title if slides and slides[0].title else os.path.basename(path)
    return DocumentIR(title=doc_title, source_path=path, slides=slides)


def ir_to_markdown(document: DocumentIR) -> str:
    lines = [f"# {document.title}", "", f"_Source: `{os.path.basename(document.source_path)}`_", ""]

    for slide in document.slides:
        heading = slide.title or f"Slide {slide.slide_number}"
        lines.append(f"## Slide {slide.slide_number}: {heading}")
        lines.append("")

        for para in slide.paragraphs:
            for line in para.splitlines():
                lines.append(line)
            lines.append("")

        for table in slide.tables:
            if table.headers:
                lines.append("| " + " | ".join(table.headers) + " |")
                lines.append("|" + "|".join(["---"] * len(table.headers)) + "|")
            for row in table.rows:
                lines.append("| " + " | ".join(row) + " |")
            lines.append("")

        for image in slide.images:
            size = ""
            if image.width_emu and image.height_emu:
                size = f" ({image.width_emu / 914400:.1f}in x {image.height_emu / 914400:.1f}in)"
            lines.append(f"_[Image: {image.shape_name}{size}]_")
            lines.append("")

        if slide.notes:
            lines.append(f"> Notes: {slide.notes}")
            lines.append("")

    return "\n".join(lines).rstrip() + "\n"


def pptx_to_markdown(path: str) -> str:
    return ir_to_markdown(pptx_to_ir(path))
