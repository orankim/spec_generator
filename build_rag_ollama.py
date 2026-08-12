"""
RAG Vector DB 구축 스크립트.

원본 데이터 형식은 Markdown(.md)을 기본으로 한다 — 사양서를 사내 PC에서 관리할 때
PPTX보다 Markdown이 다루기 쉽고(diff 가능, 사람이 직접 수정 가능), 이미
renderers/markdown_renderer.py가 표준 Markdown 포맷을 정의해 두었으므로 원본 데이터도
같은 포맷 계열로 통일하는 것이 일관적이다.

PPTX 관련 기능(parse_pptx_file)은 삭제하지 않는다 — preprocess_specs.py가 계속
사용하고, --input-dir에 PPTX가 섞여 있어도 함께 인덱싱한다. 다만 RAG 구축이 더 이상
PPTX 파일의 존재를 필수로 요구하지 않는다: --input-dir에 .md만 있어도 정상 동작한다.

임베딩 모델/서버 주소는 agent.spec_retriever.get_embeddings()를 그대로 재사용한다 —
이 스크립트(빌드)와 agent/spec_retriever.py(검색)가 서로 다른 임베딩 모델이나 Chroma
컬렉션을 쓰면 벡터 공간이 어긋나 검색이 조용히 실패하므로, 설정을 두 곳에 따로
두지 않고 한 곳(agent/spec_retriever.py)만 신뢰한다.
"""
import os

# 폐쇄망 보안 정책: 외부(HuggingFace Hub 등) 네트워크 통신 원천 차단
os.environ.setdefault("HF_HUB_OFFLINE", "1")
os.environ.setdefault("TRANSFORMERS_OFFLINE", "1")

import logging
import shutil
from glob import glob
from typing import List, Optional

from langchain_core.documents import Document
from langchain_chroma import Chroma
from langchain_text_splitters import MarkdownHeaderTextSplitter, RecursiveCharacterTextSplitter
from pptx import Presentation

logger = logging.getLogger(__name__)

# Markdown 헤딩 구조를 최대한 보존하는 chunking. H1(구분)/H2(항목) 단위로 잘라야
# "두께 측정 정밀도" 같은 개별 사양 항목이 다른 항목과 섞이지 않고 하나의 chunk가 된다
# (요청서 5절 chunking 요구사항).
_HEADERS_TO_SPLIT_ON = [("#", "header1"), ("##", "header2")]
# H2 섹션 하나가 지나치게 길면(예: 서술형 개요) 추가로 잘라 임베딩 품질을 지킨다.
_MAX_CHUNK_CHARS = 800


# ==========================================
# 1. Markdown 파일 파싱 함수 (헤딩 기반 chunking)
# ==========================================
def _extract_equipment_identity(full_text: str) -> dict:
    """
    "# 기본 정보" 아래 "## 설비명"/"## 제조사"/"## 모델명" 값을 파일 전체에서 한 번만
    뽑아, 이 파일에서 나온 모든 chunk에 공통 metadata로 붙인다. 검색 결과가 "검사
    사양 > 정확도" chunk 하나만 나와도 어떤 장비의 사양인지 알 수 있어야 하기 때문
    (Candidate Equipment 단계에서 문서 단위 그룹핑에 필요).
    """
    splitter = MarkdownHeaderTextSplitter(headers_to_split_on=_HEADERS_TO_SPLIT_ON, strip_headers=True)
    identity = {"equipment_name": None, "manufacturer": None, "model": None}
    label_to_key = {"설비명": "equipment_name", "제조사": "manufacturer", "모델명": "model"}
    for section in splitter.split_text(full_text):
        if section.metadata.get("header1") != "기본 정보":
            continue
        header2 = section.metadata.get("header2")
        key = label_to_key.get(header2)
        if key:
            identity[key] = section.page_content.strip().splitlines()[0].strip() if section.page_content.strip() else None
    return identity


def parse_markdown_file(file_path: str) -> List[Document]:
    """Markdown 사양서 1개 -> heading 기반 chunk들(Document 목록)."""
    with open(file_path, "r", encoding="utf-8") as f:
        full_text = f.read()

    file_name = os.path.basename(file_path)
    identity = _extract_equipment_identity(full_text)

    header_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=_HEADERS_TO_SPLIT_ON, strip_headers=False)
    sections = header_splitter.split_text(full_text)

    char_splitter = RecursiveCharacterTextSplitter(chunk_size=_MAX_CHUNK_CHARS, chunk_overlap=0)

    documents: List[Document] = []
    chunk_id = 0
    for section in sections:
        sub_texts = (
            char_splitter.split_text(section.page_content)
            if len(section.page_content) > _MAX_CHUNK_CHARS
            else [section.page_content]
        )
        for sub_text in sub_texts:
            if not sub_text.strip():
                continue
            metadata = {
                "source": file_path,
                "source_type": "markdown",
                "filename": file_name,
                "chunk_id": chunk_id,
                "category": section.metadata.get("header1"),
                "item": section.metadata.get("header2"),
                **{k: v for k, v in identity.items() if v},
            }
            documents.append(Document(page_content=sub_text, metadata=metadata))
            chunk_id += 1

    logger.info("[%s] Markdown 파싱 완료: %d개 chunk", file_name, len(documents))
    print(f"  └─ [{file_name}] 완료: 총 {len(documents)}개 chunk 파싱됨 (heading 기반)")
    return documents


# ==========================================
# 2. PPTX 파일 파싱 함수 (레거시 입력 지원 — 삭제하지 않음, 요청서 4절)
# ==========================================
def parse_pptx_file(file_path: str) -> list[Document]:
    prs = Presentation(file_path)
    file_name = os.path.basename(file_path)
    documents = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text_blocks = []
        table_blocks = []

        for shape in slide.shapes:
            # 텍스트 상자 추출
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_text_blocks.append(text)

            # 표(Table) 추출 -> 마크다운 형식 변환
            elif shape.has_table:
                table = shape.table
                table_str_rows = []
                for row_idx, row in enumerate(table.rows):
                    row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    row_formatted = " | ".join(row_cells)
                    table_str_rows.append(f"| {row_formatted} |")

                    if row_idx == 0:
                        table_str_rows.append("|" + "|".join(["---"] * len(row.cells)) + "|")

                if table_str_rows:
                    table_md = "\n".join(table_str_rows)
                    table_blocks.append(f"[표/규격 데이터]:\n{table_md}")

        full_slide_content = []
        if slide_text_blocks:
            full_slide_content.append("--- 슬라이드 텍스트 ---")
            full_slide_content.extend(slide_text_blocks)
        if table_blocks:
            full_slide_content.append("\n--- 사양/규격 표 데이터 ---")
            full_slide_content.extend(table_blocks)

        page_text = "\n".join(full_slide_content).strip()

        if not page_text:
            continue

        metadata = {
            "source": file_name,
            "source_type": "pptx",
            "filename": file_name,
            "chunk_id": slide_idx - 1,
            "file_path": file_path,
            "slide_number": slide_idx,
        }

        doc = Document(
            page_content=f"문서명: {file_name} (슬라이드 {slide_idx})\n\n{page_text}",
            metadata=metadata,
        )
        documents.append(doc)

    print(f"  └─ [{file_name}] 완료: 총 {len(documents)}개 슬라이드 파싱됨")
    return documents


# ==========================================
# 3. Ollama 기반 Vector DB 구축 함수
# ==========================================
def build_vector_db(
    input_dir: str,
    db_save_path: str,
    rebuild: bool = False,
    ollama_host: Optional[str] = None,
) -> Optional[Chroma]:
    """
    input_dir 안의 *.md(기본, 권장)와 *.pptx(레거시, 있으면 함께 인덱싱)를 모두 읽어
    Chroma Vector DB를 구축한다. .md만 있어도, .pptx만 있어도, 둘 다 있어도 동작한다.
    """
    from agent.spec_retriever import get_embeddings

    if rebuild and os.path.isdir(db_save_path):
        print(f"--rebuild 지정됨: 기존 Vector DB '{db_save_path}' 삭제 중...")
        shutil.rmtree(db_save_path)

    print("=== 1단계: 사양서 파일 파싱 시작 ===")
    md_files = sorted(glob(os.path.join(input_dir, "*.md")))
    pptx_files = sorted(glob(os.path.join(input_dir, "*.pptx")))

    if not md_files and not pptx_files:
        print(f"경고: '{input_dir}' 경로에 .md 또는 .pptx 파일이 없습니다.")
        return None

    all_documents: List[Document] = []
    for md_file in md_files:
        all_documents.extend(parse_markdown_file(md_file))
    for pptx_file in pptx_files:
        all_documents.extend(parse_pptx_file(pptx_file))

    print(f"\n파일 {len(md_files)}개(.md) + {len(pptx_files)}개(.pptx)에서 총 {len(all_documents)}개의 chunk가 준비되었습니다.")

    print("\n=== 2단계: 사내 Ollama 임베딩 엔진 연결 ===")
    embeddings = get_embeddings(ollama_host)

    print("\n=== 3단계: Chroma Vector DB 생성 및 저장 ===")
    vector_store = Chroma.from_documents(
        documents=all_documents,
        embedding=embeddings,
        persist_directory=db_save_path,
    )
    print(f"성공! Vector DB가 '{db_save_path}' 경로에 저장되었습니다.\n")
    return vector_store


# 이전 이름과의 호환 — 다른 스크립트/문서가 이 이름으로 호출하고 있을 수 있다.
def build_vector_db_with_ollama(pptx_folder_path: str, db_save_path: str) -> Optional[Chroma]:
    return build_vector_db(pptx_folder_path, db_save_path)


# ==========================================
# 4. 검색 테스트 함수
# ==========================================
def test_search(query: str, db_save_path: str, ollama_host: Optional[str] = None) -> None:
    from agent.spec_retriever import get_embeddings

    print(f"\n=== 검색 테스트: '{query}' ===")
    embeddings = get_embeddings(ollama_host)
    vector_store = Chroma(persist_directory=db_save_path, embedding_function=embeddings)

    results = vector_store.similarity_search(query, k=3)
    for i, doc in enumerate(results, start=1):
        source = doc.metadata.get("filename", doc.metadata.get("source", "?"))
        print(f"\n[검색 결과 {i}] (출처: {source})")
        print("-" * 50)
        print(doc.page_content)
        print("-" * 50)


# ==========================================
# 실행부
# ==========================================
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Markdown(.md) 사양서로 RAG Vector DB 구축 (.pptx도 함께 지원)")
    parser.add_argument("--input-dir", default="./sample_specs", help="파싱할 사양서 폴더 (.md 우선, .pptx도 지원)")
    parser.add_argument(
        "--pptx-folder", default=None,
        help="[레거시 별칭] --input-dir와 동일하게 동작한다. 이름과 달리 .md/.pptx를 모두 스캔한다.",
    )
    parser.add_argument("--db-path", default="./chroma_db_specs", help="저장할 Vector DB 폴더")
    parser.add_argument("--rebuild", action="store_true", help="기존 Vector DB를 삭제하고 새로 구축한다")
    args = parser.parse_args()

    input_dir = args.pptx_folder if args.pptx_folder is not None else args.input_dir
    os.makedirs(input_dir, exist_ok=True)

    build_vector_db(input_dir, args.db_path, rebuild=args.rebuild)

    # 검색 테스트가 필요하면 주석을 해제하세요.
    # test_search("두께 측정 정밀도", args.db_path)
