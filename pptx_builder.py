import os
from copy import deepcopy
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PPTXBuilder:
    def __init__(self, template_path: str = "template.pptx"):
        """
        PPTX 템플릿 파일 경로를 받아 초기화합니다.
        """
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"템플릿 파일 '{template_path}'을(를) 찾을 수 없습니다. 마스터 템플릿을 준비해 주세요.")
        
        self.template_path = template_path

    def _replace_text_in_shape(self, shape, replacements: Dict[str, str]):
        """
        도형/텍스트 상자 안의 태그({{TAG}})를 실제 데이터로 치환합니다.
        """
        if not shape.has_text_frame:
            return

        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for tag, value in replacements.items():
                if tag in paragraph.text:
                    # 기존 텍스트 치환
                    paragraph.text = paragraph.text.replace(tag, str(value))
                    # 폰트 스타일 유지 (필요시 조정 가능)
                    paragraph.font.size = Pt(14)
                    paragraph.font.name = "맑은 고딕"

    def _add_table_row(self, table):
        """
        python-pptx는 기존 표에 행을 추가하는 공개 API(add_row)를 제공하지 않으므로,
        마지막 행의 XML(<a:tr>)을 복제하여 서식을 유지한 채 행을 추가합니다.
        """
        new_tr = deepcopy(table._tbl.tr_lst[-1])
        table._tbl.append(new_tr)

    def _populate_spec_table(self, slide, spec_items: List[Dict[str, str]]):
        """
        슬라이드에 있는 기존 표(Table)를 찾아 사양 데이터를 채워 넣습니다.
        """
        table_shape = None
        for shape in slide.shapes:
            if shape.has_table:
                table_shape = shape
                break

        if not table_shape:
            print("⚠️ 경고: 해당 슬라이드에서 표(Table)를 찾지 못했습니다.")
            return

        table = table_shape.table

        # 헤더 행을 제외한 기존 데이터 행 제거 또는 새로 추가
        # python-pptx는 행 삭제가 까다로우므로, 부족한 행만큼 새로 추가하는 방식 적용
        needed_rows = len(spec_items) + 1  # 데이터 행 + 헤더 1행
        current_rows = len(table.rows)

        # 행이 부족하면 추가
        while len(table.rows) < needed_rows:
            self._add_table_row(table)

        # 데이터 채우기 (row_idx = 0 은 헤더)
        for idx, item in enumerate(spec_items, start=1):
            row_cells = table.rows[idx].cells
            
            # 각 셀에 데이터 바인딩
            row_cells[0].text = item.get("category", "")
            row_cells[1].text = item.get("item", "")
            row_cells[2].text = item.get("value", "")
            row_cells[3].text = item.get("note", "")

            # 서식 적용 (폰트 크기, 정렬 등)
            for cell_idx, cell in enumerate(row_cells):
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = "맑은 고딕"
                    paragraph.font.size = Pt(11)
                    if cell_idx == 0:  # 카테고리는 중앙 정렬
                        paragraph.alignment = PP_ALIGN.CENTER

    def build(self, spec_data: Dict[str, Any], output_path: str = "generated_spec.pptx") -> str:
        """
        JSON 데이터를 받아 PPTX 사양서를 합성 및 저장합니다.
        """
        print(f"\n=== PPTX 사양서 생성 시작: '{output_path}' ===")
        prs = Presentation(self.template_path)

        # 태그 치환용 데이터 맵핑
        replacements = {
            "{{EQUIPMENT_NAME}}": spec_data.get("equipment_name", "설비 사양서"),
            "{{OVERVIEW}}": spec_data.get("overview", "설비 개요 정보가 없습니다."),
            "{{TARGET_CAPACITY}}": spec_data.get("target_capacity", "사양 정보 없음")
        }

        # 1. 모든 슬라이드를 순회하며 텍스트 태그 치환
        for slide in prs.slides:
            for shape in slide.shapes:
                self._replace_text_in_shape(shape, replacements)

        # 2. 두 번째 슬라이드(또는 표가 있는 슬라이드)에 상세 사양 표 채우기
        spec_table_data = spec_data.get("spec_table", [])
        if len(prs.slides) > 1 and spec_table_data:
            # 2번째 슬라이드에 표 데이터 채우기
            self._populate_spec_table(prs.slides[1], spec_table_data)

        # 3. 파일 저장
        prs.save(output_path)
        print(f" 성공! 최종 사양서가 저장되었습니다: {os.path.abspath(output_path)}")
        return output_path


# ==========================================
# 단독 테스트 연동
# ==========================================
if __name__ == "__main__":
    # generator.py에서 생성된 형태의 샘플 JSON 데이터
    sample_spec_json = {
        "equipment_name": "300mm 고진공 플라즈마 식각 설비",
        "overview": "본 설비는 300mm 웨이퍼 표면의 미세 패턴을 고진공 환경에서 식각하기 위한 전용 장비입니다.",
        "target_capacity": "시간당 30장 (30 wph)",
        "spec_table": [
            {"category": "전기/전력", "item": "정격 전압", "value": "3Phi 380V 60Hz", "note": "전압 변동률 ±5% 이내"},
            {"category": "진공 사양", "item": "도달 진공도", "value": "1.0 x 10^-6 Torr", "note": "터보분자펌프 적용"},
            {"category": "치수/중량", "item": "설비 크기", "value": "2100(W) x 1800(D) x 2000(H) mm", "note": "유지보수 공간 제외"},
            {"category": "제어 방식", "item": "Main PLC", "value": "Siemens S7-1500", "note": "사내 Ethernet 통신 규격"}
        ]
    }

    # 테스트 실행
    try:
        builder = PPTXBuilder(template_path="template.pptx")
        output_file = builder.build(sample_spec_json, output_path="test_spec_output.pptx")
    except FileNotFoundError as e:
        print(e)
        print("💡 TIP: 파워포인트를 열어 'template.pptx' 파일을 프로젝트 폴더에 저장한 후 다시 실행해 주세요.")
