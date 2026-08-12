"""
전극 검사기 사양서 전용 9섹션 마스터 템플릿(template_electrode.pptx) 생성기.
기존 make_template.py(2슬라이드, 범용 사양서)는 건드리지 않고 별도 파일로 생성한다.

각 슬라이드 제목 도형에 {{SECTION:...}} 마커를 심어두어, 빌더가 슬라이드
순서가 아니라 마커로 슬라이드를 찾을 수 있게 한다 (슬라이드 순서가 바뀌어도
안전하게 동작).
"""
from pptx import Presentation
from pptx.util import Inches, Pt

SECTIONS = [
    ("COVER", "표지"),
    ("GENERAL", "1. General Specification (장비 개요)"),
    ("INSPECTION_TARGET", "2. Inspection Target (검사 대상)"),
    ("MEASUREMENT_PERFORMANCE", "3. Measurement Performance (측정 성능)"),
    ("INSPECTION_PERFORMANCE", "4. Inspection Performance (검사 성능 / 결함 검출)"),
    ("SYSTEM_CONFIG", "5. System Configuration (시스템 구성)"),
    ("INTERFACE", "6. Interface (연동 인터페이스)"),
    ("ENVIRONMENT", "7. Environment & Safety (설치 환경 / 안전)"),
    ("NOTES", "8. Notes & Evidence (비고 / 근거 자료)"),
]

TABLE_COLUMN_HEADERS = ["항목", "값", "단위", "근거(출처)"]


def _add_title_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    marker = slide.shapes.add_textbox(Inches(0.1), Inches(0.1), Inches(0.1), Inches(0.1))
    marker.text_frame.text = "{{SECTION:COVER}}"
    marker.text_frame.paragraphs[0].font.size = Pt(1)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "{{EQUIPMENT_NAME}}"
    p.font.size = Pt(32)
    p.font.bold = True

    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(2))
    tf = sub_box.text_frame
    tf.paragraphs[0].text = "검사 대상: {{MATERIAL}}"
    tf.paragraphs[0].font.size = Pt(16)
    p2 = tf.add_paragraph()
    p2.text = "측정 원리: {{MEASUREMENT_PRINCIPLE}}"
    p2.font.size = Pt(16)
    return slide


def _add_table_slide(prs, section_key: str, title: str, n_data_rows: int = 6, table_cols=None):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    marker = slide.shapes.add_textbox(Inches(0.1), Inches(0.05), Inches(0.1), Inches(0.1))
    marker.text_frame.text = f"{{{{SECTION:{section_key}}}}}"
    marker.text_frame.paragraphs[0].font.size = Pt(1)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(22)
    tp.font.bold = True

    cols = table_cols or TABLE_COLUMN_HEADERS
    n_rows = n_data_rows + 1
    table_shape = slide.shapes.add_table(n_rows, len(cols), Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    table = table_shape.table
    for c_idx, header in enumerate(cols):
        table.rows[0].cells[c_idx].text = header
    return slide


def build_electrode_template(output_path: str = "template_electrode.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs)
    _add_table_slide(prs, "GENERAL", SECTIONS[1][1], n_data_rows=4, table_cols=["항목", "값"])
    _add_table_slide(prs, "INSPECTION_TARGET", SECTIONS[2][1], n_data_rows=7, table_cols=["항목", "값"])
    _add_table_slide(prs, "MEASUREMENT_PERFORMANCE", SECTIONS[3][1], n_data_rows=10)
    _add_table_slide(prs, "INSPECTION_PERFORMANCE", SECTIONS[4][1], n_data_rows=10)
    _add_table_slide(prs, "SYSTEM_CONFIG", SECTIONS[5][1], n_data_rows=8, table_cols=["항목", "값"])
    _add_table_slide(prs, "INTERFACE", SECTIONS[6][1], n_data_rows=6, table_cols=["항목", "값"])
    _add_table_slide(prs, "ENVIRONMENT", SECTIONS[7][1], n_data_rows=8, table_cols=["항목", "값"])
    _add_table_slide(prs, "NOTES", SECTIONS[8][1], n_data_rows=6, table_cols=["구분", "내용"])

    prs.save(output_path)
    print(f"성공! '{output_path}' 파일이 생성되었습니다. (총 {len(prs.slides)}개 슬라이드)")


if __name__ == "__main__":
    build_electrode_template()
