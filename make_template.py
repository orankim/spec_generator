# make_template.py
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# 슬라이드 1 (표지/개요)
blank_slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_slide_layout)

# 제목 텍스트 박스
txBox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "{{EQUIPMENT_NAME}}"
p.font.size = Pt(28)
p.font.bold = True

# 개요 및 용량 텍스트 박스
txBox2 = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(3.5))
tf2 = txBox2.text_frame

p1 = tf2.paragraphs[0]
p1.text = "[개요]\n{{OVERVIEW}}"
p1.font.size = Pt(14)

p2 = tf2.add_paragraph()
p2.text = "\n[처리 용량/성능]\n{{TARGET_CAPACITY}}"
p2.font.size = Pt(14)

# 슬라이드 2 (상세 사양 표)
slide2 = prs.slides.add_slide(blank_slide_layout)

# 제목
title_box = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_tf = title_box.text_frame
title_p = title_tf.paragraphs[0]
title_p.text = "상세 기술 사양서"
title_p.font.size = Pt(24)
title_p.font.bold = True

# 표 생성 (2행 4열 샘플)
table_shape = slide2.shapes.add_table(2, 4, Inches(1), Inches(1.8), Inches(8), Inches(4.5))
table = table_shape.table

# 표 헤더 세팅
headers = ["구분", "항목", "사양값", "비고"]
for idx, header in enumerate(headers):
    table.rows[0].cells[idx].text = header

# 저장
prs.save("template.pptx")
print("성공! 'template.pptx' 파일이 생성되었습니다.")
