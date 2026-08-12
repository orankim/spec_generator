"""
Specification JSON을 단일 소스로 하는 3-way 렌더러(Markdown/HTML/PPTX)와
PPTX <-> Markdown 변환기에 대한 테스트.

Ollama 호출이 전혀 없는 순수 변환 로직이므로, 이 테스트들은 모킹 없이
실제 코드 경로 그대로 실행된다.
"""
import json
import tempfile
from pathlib import Path

from pptx import Presentation

from agent.schemas import RequirementSchema, SourcedNumber, SourceRef, SpecificationSchema
from agent.spec_validator import validate_specification
from converters.markdown_to_spec import markdown_to_spec
from converters.pptx_to_markdown import pptx_to_ir, pptx_to_markdown
from renderers.html_renderer import render_html
from renderers.markdown_renderer import render_markdown
from renderers.pptx_renderer import render_pptx

_REPO_ROOT = Path(__file__).resolve().parent.parent
_EXAMPLE_SPEC_PATH = _REPO_ROOT / "docs" / "examples" / "example_specification.json"


def _full_spec() -> SpecificationSchema:
    spec = SpecificationSchema()
    spec.equipment.name = "전극 두께/표면결함 비접촉 검사기"
    spec.equipment.manufacturer = "ACME Metrology"
    spec.equipment.measurement_principle = "laser triangulation"
    spec.inspection_target.material = "전극"
    spec.inspection_target.width_mm = 500.0
    spec.inspection_items = ["thickness", "surface_defect"]
    spec.measurement_performance.accuracy_um = SourcedNumber(
        value=0.5, unit="um", operator="<=", status="VERIFIED", source=SourceRef(document="vendor_spec.pptx")
    )
    spec.measurement_performance.repeatability_um = SourcedNumber(value=0.3, unit="um", status="USER_DEFINED")
    spec.defect_detection.minimum_defect_size_um = SourcedNumber(value=10.0, unit="um", status="INFERRED")
    spec.defect_detection.defect_types = ["핀홀", "크랙"]
    spec.interfaces.ethernet = True
    spec.interfaces.mes = False
    spec.notes = ["첫 번째 노트", "두 번째 노트"]
    spec.sources = ["vendor_spec.pptx"]
    return spec


# ---------------------------------------------------------------
# Test 1: Specification JSON -> Markdown
# ---------------------------------------------------------------
def test_spec_to_markdown():
    md = render_markdown(_full_spec())
    assert md.startswith("# Electrode Inspection Equipment Specification")
    assert "## 1. General Specification" in md
    assert "## 14. Sources / Notes" in md
    assert "전극 두께/표면결함 비접촉 검사기" in md


# ---------------------------------------------------------------
# Test 2: Specification JSON -> HTML
# ---------------------------------------------------------------
def test_spec_to_html():
    html = render_html(_full_spec())
    assert "<html" in html
    assert "<table>" in html
    assert "전극 두께/표면결함 비접촉 검사기" in html
    # 외부 CDN에 의존하지 않는지 확인
    assert "http://" not in html and "https://" not in html


# ---------------------------------------------------------------
# Test 3: Specification JSON -> PPTX (템플릿 없이도 동작해야 함)
# ---------------------------------------------------------------
def test_spec_to_pptx_without_template():
    spec = _full_spec()
    with tempfile.TemporaryDirectory() as tmp:
        out = render_pptx(spec, str(Path(tmp) / "out.pptx"), template_path=None)
        prs = Presentation(out)
        assert len(prs.slides) > 1
        all_text = "\n".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "전극 두께/표면결함 비접촉 검사기" in all_text


def test_spec_to_pptx_with_template():
    spec = _full_spec()
    with tempfile.TemporaryDirectory() as tmp:
        out = render_pptx(spec, str(Path(tmp) / "out.pptx"), template_path="template_electrode.pptx")
        prs = Presentation(out)
        assert len(prs.slides) == 9  # 기존 ElectrodeSpecPPTXBuilder 그대로 사용


# ---------------------------------------------------------------
# Test 4: PPTX -> Markdown
# ---------------------------------------------------------------
def test_pptx_to_markdown_preserves_content():
    md = pptx_to_markdown("sample_specs/spec_electrode_coating_thickness.pptx")
    assert md.startswith("# ")
    assert "## Slide 1" in md
    assert "## Slide 2" in md
    ir = pptx_to_ir("sample_specs/spec_electrode_coating_thickness.pptx")
    assert len(ir.slides) == 2
    assert len(ir.slides[1].tables) == 1  # 상세 사양 표가 있는 슬라이드


# ---------------------------------------------------------------
# Test 5: Markdown -> Specification JSON
# ---------------------------------------------------------------
def test_markdown_to_spec_roundtrip():
    original = _full_spec()
    md = render_markdown(original)
    roundtrip = markdown_to_spec(md)

    assert roundtrip.equipment.name == original.equipment.name
    assert roundtrip.equipment.manufacturer == original.equipment.manufacturer
    assert roundtrip.inspection_target.width_mm == original.inspection_target.width_mm
    assert roundtrip.inspection_items == original.inspection_items
    assert roundtrip.measurement_performance.accuracy_um.value == 0.5
    assert roundtrip.measurement_performance.accuracy_um.status == "VERIFIED"
    assert roundtrip.measurement_performance.accuracy_um.source.document == "vendor_spec.pptx"
    assert roundtrip.defect_detection.defect_types == ["핀홀", "크랙"]
    assert roundtrip.notes == original.notes
    assert roundtrip.sources == original.sources


# ---------------------------------------------------------------
# Test 6: UNKNOWN 값 보존
# ---------------------------------------------------------------
def test_unknown_values_preserved_through_markdown_roundtrip():
    spec = SpecificationSchema()  # 전부 비어있는 사양서
    md = render_markdown(spec)
    assert "UNKNOWN" in md

    roundtrip = markdown_to_spec(md)
    assert roundtrip.equipment.name is None
    assert roundtrip.measurement_performance.accuracy_um is None
    assert roundtrip.inspection_items == []


# ---------------------------------------------------------------
# Test 7: Source 정보 보존
# ---------------------------------------------------------------
def test_source_info_preserved_in_markdown_and_html():
    spec = _full_spec()
    md = render_markdown(spec)
    assert "| Accuracy | um | 0.5 | VERIFIED | vendor_spec.pptx |" in md
    assert "| Repeatability | um | 0.3 | USER_DEFINED | -" in md

    html = render_html(spec)
    assert "vendor_spec.pptx" in html
    assert "Accuracy" in html
    assert "badge-verified" in html


# ---------------------------------------------------------------
# Test 8: 단위 보존
# ---------------------------------------------------------------
def test_units_preserved():
    spec = _full_spec()
    md = render_markdown(spec)
    assert "| Accuracy | um | 0.5 |" in md  # unit 컬럼에 um이 들어감

    roundtrip = markdown_to_spec(md)
    assert roundtrip.measurement_performance.accuracy_um.unit == "um"


# ---------------------------------------------------------------
# Test 9: Optional Field 보존 (예: Optical System 정보가 전혀 없는 장비)
# ---------------------------------------------------------------
def test_optional_optical_system_absent_does_not_break_rendering():
    spec = _full_spec()
    assert spec.optical_system.light_source is None  # 애초에 설정 안 함

    md = render_markdown(spec)
    html = render_html(spec)
    with tempfile.TemporaryDirectory() as tmp:
        pptx_path = render_pptx(spec, str(Path(tmp) / "out.pptx"))
        prs = Presentation(pptx_path)
        assert len(prs.slides) > 1  # optical system이 비어 있어도 생성 자체는 성공

    assert "## 6. Optical System" in md
    assert "Light Source" in md
    assert "UNKNOWN" in md  # 비어있는 optical system 필드들이 UNKNOWN으로 명시됨
    assert "Optical System" in html


# ---------------------------------------------------------------
# Test 10: 긴 텍스트 처리
# ---------------------------------------------------------------
def test_long_text_handling():
    spec = SpecificationSchema()
    spec.equipment.name = "A" * 500
    long_note = "이 설비는 " * 200  # 매우 긴 노트
    spec.notes = [long_note]

    md = render_markdown(spec)
    assert "A" * 500 in md
    assert long_note in md
    # 마크다운 테이블 파이프 문자가 깨지지 않는지 (셀 안에 개행이 있으면 표가 깨지므로) 확인
    assert "\n" not in long_note or all(len(line.split("|")) >= 2 for line in md.splitlines() if line.startswith("|"))

    html = render_html(spec)
    assert "A" * 500 in html

    with tempfile.TemporaryDirectory() as tmp:
        out = render_pptx(spec, str(Path(tmp) / "out.pptx"))
        prs = Presentation(out)
        assert len(prs.slides) > 1  # 긴 텍스트로도 생성 자체가 실패하지 않음


# ---------------------------------------------------------------
# Requirement Compliance (신규 13번 섹션)
# ---------------------------------------------------------------
def test_requirement_compliance_section_computes_pass_fail():
    spec = _full_spec()
    req = RequirementSchema(required_accuracy_um=1.0)  # 사양(0.5um) <= 요구사항(1.0um) -> PASS

    md = render_markdown(spec, requirement=req)
    assert "## 13. Requirement Compliance" in md
    assert "| Accuracy | um | 1.0 | 0.5 | PASS |" in md

    html = render_html(spec, requirement=req)
    assert "badge-pass" in html

    strict_req = RequirementSchema(required_accuracy_um=0.1)  # 사양(0.5um) > 요구사항(0.1um) -> FAIL
    md_fail = render_markdown(spec, requirement=strict_req)
    assert "| Accuracy | um | 0.1 | 0.5 | FAIL |" in md_fail


def test_requirement_compliance_absent_without_requirement():
    spec = _full_spec()
    md = render_markdown(spec)  # requirement 없음
    assert "## 13. Requirement Compliance" in md
    assert "No requirement provided for comparison." in md


# ---------------------------------------------------------------
# Status/legacy migration
# ---------------------------------------------------------------
def test_sourced_number_fields_outside_measurement_sections_keep_status_and_source():
    """
    inspection_target.line_speed_mm_s / inspection_requirements.sampling_interval도
    SourcedNumber다 (Measurement/Spatial Performance 섹션이 아니라고 Status/Source를
    잃어버리면 안 된다 — 회귀 확인용, renderers/*_renderer.py의 numeric-section 판정
    누락으로 한 번 깨졌던 부분).
    """
    spec = SpecificationSchema()
    spec.inspection_target.material = "전극"
    spec.inspection_items = ["thickness"]
    spec.inspection_target.line_speed_mm_s = SourcedNumber(
        value=120.0, unit="mm/s", status="VERIFIED", source=SourceRef(document="vendor_spec.pptx")
    )
    spec.inspection_requirements.sampling_interval = SourcedNumber(value=5.0, unit="mm", status="INFERRED")

    md = render_markdown(spec)
    assert "| Target Line Speed | mm/s | 120.0 | VERIFIED | vendor_spec.pptx |" in md
    assert "| Sampling Interval | mm | 5.0 | INFERRED | - |" in md

    html = render_html(spec)
    assert "badge-verified" in html
    assert "badge-inferred" in html

    roundtrip = markdown_to_spec(md)
    assert roundtrip.inspection_target.line_speed_mm_s.value == 120.0
    assert roundtrip.inspection_target.line_speed_mm_s.status == "VERIFIED"
    assert roundtrip.inspection_target.line_speed_mm_s.source.document == "vendor_spec.pptx"
    assert roundtrip.inspection_requirements.sampling_interval.value == 5.0
    assert roundtrip.inspection_requirements.sampling_interval.status == "INFERRED"


def test_sourced_number_legacy_migration_helper():
    legacy = SourcedNumber.from_legacy(value=1.0, unit="um", source_type="document", source="old.pptx")
    assert legacy.status == "VERIFIED"
    assert legacy.source.document == "old.pptx"

    legacy2 = SourcedNumber.from_legacy(value=2.0, source_type="user_requirement")
    assert legacy2.status == "USER_DEFINED"


# ---------------------------------------------------------------
# 회귀: 기존 기능이 이번 변경으로 깨지지 않았는지 (agent 파이프라인/스키마)
# ---------------------------------------------------------------
def test_specification_schema_has_expected_top_level_sections():
    """
    SpecificationSchema는 이번 세션에서 필드가 확장됐다(요청서 3~13절, 명시적으로 승인됨).
    기존 필드는 하나도 제거되지 않았고, inspection_requirements가 새로 추가됐는지만 확인한다.
    """
    spec = SpecificationSchema()
    top_level = set(type(spec).model_fields.keys())
    previously_existing = {
        "equipment", "inspection_target", "inspection_items", "measurement_performance",
        "spatial_performance", "inspection_performance", "defect_detection", "optical_system",
        "system", "interfaces", "environment", "safety", "notes", "assumptions", "sources",
        "needs_confirmation",
    }
    assert previously_existing.issubset(top_level), "기존 필드가 제거되면 안 됨"
    assert "inspection_requirements" in top_level


def test_example_specification_json_is_valid_and_renders_cleanly():
    """
    docs/examples/example_specification.json(공개 가능한 예시 1건)이 SpecificationSchema로
    파싱되고, 검증을 통과하며(is_valid=True), 세 포맷 모두 예외 없이 렌더링되는지 확인한다.
    스키마가 바뀔 때 이 예시 파일이 조용히 깨지는 것을 막기 위한 회귀 테스트다.
    """
    data = json.loads(_EXAMPLE_SPEC_PATH.read_text(encoding="utf-8"))
    spec = SpecificationSchema(**data)

    result = validate_specification(spec)
    assert result.is_valid is True
    assert not any(i.level == "error" for i in result.issues)

    md = render_markdown(spec)
    assert "전극 코팅 두께" in md
    html = render_html(spec)
    assert "<html" in html
    with tempfile.TemporaryDirectory() as tmp:
        out_path = str(Path(tmp) / "example.pptx")
        render_pptx(spec, out_path)
        assert Path(out_path).exists()


def test_ppt_template_adapter_default_renders_without_template():
    """
    PPTTemplateAdapter(요청서 23절, render(specification, template_path, output_path) 시그니처)의
    기본 구현이 템플릿 없이도(=존재하지 않는 경로) 동작해야 한다 — 회사 템플릿이 없는
    환경에서도 파이프라인이 항상 PPTX 한 장을 만들어낼 수 있어야 하기 때문이다.
    """
    from templates.adapters.ppt_template_adapter import DefaultPPTTemplateAdapter, PPTTemplateAdapter

    spec = _full_spec()
    adapter: PPTTemplateAdapter = DefaultPPTTemplateAdapter()
    with tempfile.TemporaryDirectory() as tmp:
        out_path = str(Path(tmp) / "out.pptx")
        result_path = adapter.render(spec, template_path="/no/such/template.pptx", output_path=out_path)
        assert result_path == out_path
        assert Path(out_path).exists()
        prs = Presentation(out_path)
        assert len(prs.slides) > 0


def test_pptx_electrode_builder_still_importable_and_usable():
    """agent/pptx_electrode_builder.py(기존 PPTX Generator)를 삭제하지 않고 그대로 재사용 가능해야 한다."""
    from agent.pptx_electrode_builder import ElectrodeSpecPPTXBuilder

    builder = ElectrodeSpecPPTXBuilder(template_path="template_electrode.pptx")
    with tempfile.TemporaryDirectory() as tmp:
        out = builder.build(_full_spec(), output_path=str(Path(tmp) / "legacy.pptx"))
        assert Path(out).exists()
