"""
ElectrodeSpecPPTXBuilder — SpecificationSchema를 9섹션 template_electrode.pptx에
채워 넣는다. LLM은 PPTX를 생성하지 않고, 이 모듈(순수 Python)이 담당한다
(기획안 14절).

기존 pptx_builder.PPTXBuilder를 상속해 태그 치환(_replace_text_in_shape)과
표 행 추가(_add_table_row, python-pptx에 공개 API가 없어 XML 복제로 구현된
저수준 유틸리티)를 그대로 재사용한다.
"""
from __future__ import annotations

import os
import sys
from typing import List, Optional

from pptx import Presentation
from pptx.util import Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from pptx_builder import PPTXBuilder  # noqa: E402  (기존 프로젝트 최상위 모듈 재사용)

from .schemas import SourcedNumber, SpecificationSchema  # noqa: E402


def _num_row(label: str, sn: Optional[SourcedNumber]) -> List[str]:
    if sn is None or sn.value is None:
        return [label, "N/A", "", ""]
    if sn.source and sn.source.document:
        source = sn.source.document
    else:
        source = {"USER_DEFINED": "사용자 요구사항", "INFERRED": "AI 추정", "UNKNOWN": ""}.get(sn.status, "")
    return [label, str(sn.value), sn.unit or "", source]


def _plain_row(label: str, value) -> List[str]:
    if value is None or value == "" or value == []:
        return [label, "N/A"]
    if isinstance(value, list):
        return [label, ", ".join(str(v) for v in value)]
    if isinstance(value, bool):
        return [label, "지원" if value else "미지원"]
    return [label, str(value)]


class ElectrodeSpecPPTXBuilder(PPTXBuilder):
    def _find_slide_by_section(self, prs, section_key: str):
        marker = f"{{{{SECTION:{section_key}}}}}"
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and marker in shape.text_frame.text:
                    return slide
        return None

    def _fill_table(self, slide, rows: List[List[str]]):
        table_shape = next((s for s in slide.shapes if s.has_table), None)
        if not table_shape:
            print("⚠️ 경고: 섹션 슬라이드에서 표를 찾지 못했습니다.")
            return
        table = table_shape.table

        needed_rows = len(rows) + 1  # 헤더 + 데이터
        while len(table.rows) < needed_rows:
            self._add_table_row(table)

        for r_idx, row_values in enumerate(rows, start=1):
            cells = table.rows[r_idx].cells
            for c_idx in range(len(cells)):
                text = str(row_values[c_idx]) if c_idx < len(row_values) else ""
                cells[c_idx].text = text
                for p in cells[c_idx].text_frame.paragraphs:
                    p.font.name = "맑은 고딕"
                    p.font.size = Pt(11)

    def build(self, spec: SpecificationSchema, output_path: str) -> str:
        print(f"\n=== 전극 검사기 사양서 PPTX 생성 시작: '{output_path}' ===")
        prs = Presentation(self.template_path)

        equipment_name = spec.equipment.name or f"{spec.inspection_target.material or '전극'} 검사기 사양서"
        replacements = {
            "{{EQUIPMENT_NAME}}": equipment_name,
            "{{MATERIAL}}": spec.inspection_target.material or "N/A",
            "{{MEASUREMENT_PRINCIPLE}}": spec.equipment.measurement_principle or "N/A",
        }
        for slide in prs.slides:
            for shape in slide.shapes:
                self._replace_text_in_shape(shape, replacements)

        eq = spec.equipment
        general_slide = self._find_slide_by_section(prs, "GENERAL")
        if general_slide:
            self._fill_table(general_slide, [
                _plain_row("설비명", eq.name),
                _plain_row("제조사", eq.manufacturer),
                _plain_row("모델명", eq.model),
                _plain_row("측정 원리", eq.measurement_principle),
            ])

        it = spec.inspection_target
        target_slide = self._find_slide_by_section(prs, "INSPECTION_TARGET")
        if target_slide:
            self._fill_table(target_slide, [
                _plain_row("검사 대상(material)", it.material),
                _plain_row("제품 유형", it.product_type),
                _plain_row("폭 (mm)", it.width_mm),
                _plain_row("길이 (mm)", it.length_mm),
                _plain_row("두께 (um)", it.thickness_um),
                _plain_row("기판(substrate)", it.substrate),
                _plain_row("검사 방향", it.inspection_direction),
            ])

        mp, sp = spec.measurement_performance, spec.spatial_performance
        measurement_slide = self._find_slide_by_section(prs, "MEASUREMENT_PERFORMANCE")
        if measurement_slide:
            self._fill_table(measurement_slide, [
                _num_row("측정 범위", mp.measurement_range),
                _num_row("분해능 (Resolution)", mp.resolution_um),
                _num_row("정확도 (Accuracy)", mp.accuracy_um),
                _num_row("반복성 (Repeatability)", mp.repeatability_um),
                _num_row("재현성 (Reproducibility)", mp.reproducibility_um),
                _num_row("FOV", sp.fov_mm),
                _num_row("X 분해능", sp.x_resolution_um),
                _num_row("Y 분해능", sp.y_resolution_um),
                _num_row("Z 분해능", sp.z_resolution_um),
                _num_row("샘플링 간격", sp.sampling_interval_um),
            ])

        ip, dd = spec.inspection_performance, spec.defect_detection
        inspection_slide = self._find_slide_by_section(prs, "INSPECTION_PERFORMANCE")
        if inspection_slide:
            self._fill_table(inspection_slide, [
                _num_row("스캔 속도", ip.scan_speed_mm_s),
                _num_row("라인 속도", ip.line_speed_mm_s),
                _num_row("측정 속도", ip.measurement_speed),
                _num_row("Tact Time", ip.tact_time_s),
                _num_row("검사 폭", ip.inspection_width_mm),
                _num_row("최소 검출 결함 크기", dd.minimum_defect_size_um),
                _plain_row("검출 가능 결함 유형", dd.defect_types),
                _num_row("결함 검출 정확도", dd.defect_detection_accuracy),
                _num_row("오검출률 (False Positive)", dd.false_positive_rate),
                _num_row("미검출률 (False Negative)", dd.false_negative_rate),
            ])

        sysc, opt = spec.system, spec.optical_system
        system_slide = self._find_slide_by_section(prs, "SYSTEM_CONFIG")
        if system_slide:
            self._fill_table(system_slide, [
                _plain_row("자동화 수준", sysc.automation_level),
                _plain_row("스테이지", sysc.stage),
                _plain_row("구동계", sysc.motion_system),
                _plain_row("컨트롤러", sysc.controller),
                _plain_row("소프트웨어", sysc.software),
                _plain_row("데이터 출력", sysc.data_output),
                _plain_row("광원 / 파장", f"{opt.light_source or ''} {opt.wavelength or ''}".strip() or None),
                _plain_row("광학 방식 / 센서", f"{opt.optical_method or ''} {opt.sensor_type or ''}".strip() or None),
            ])

        iface = spec.interfaces
        interface_slide = self._find_slide_by_section(prs, "INTERFACE")
        if interface_slide:
            self._fill_table(interface_slide, [
                _plain_row("Ethernet", iface.ethernet),
                _plain_row("Digital I/O", iface.digital_io),
                _plain_row("PLC 연동", iface.plc),
                _plain_row("MES 연동", iface.mes),
                _plain_row("OPC-UA", iface.opc_ua),
                _plain_row("기타 인터페이스", iface.other_interfaces),
            ])

        env, safety = spec.environment, spec.safety
        env_slide = self._find_slide_by_section(prs, "ENVIRONMENT")
        if env_slide:
            self._fill_table(env_slide, [
                _plain_row("동작 온도", env.operating_temperature),
                _plain_row("습도", env.humidity),
                _plain_row("설치 공간", env.installation_space),
                _plain_row("전원", env.power),
                _plain_row("진동 요구사항", env.vibration_requirement),
                _plain_row("안전 규격", safety.safety_standard),
                _plain_row("인터록", safety.interlock),
                _plain_row("비상정지", safety.emergency_stop),
            ])

        notes_slide = self._find_slide_by_section(prs, "NOTES")
        if notes_slide:
            rows = []
            for note in spec.notes:
                rows.append(_plain_row("Note", note))
            for assumption in spec.assumptions:
                rows.append(_plain_row("가정(Assumption)", assumption))
            if spec.needs_confirmation:
                rows.append(_plain_row("확인 필요 항목", spec.needs_confirmation))
            if spec.sources:
                rows.append(_plain_row("참고 문서(Sources)", spec.sources))
            if not rows:
                rows.append(["-", "-"])
            self._fill_table(notes_slide, rows)

        prs.save(output_path)
        print(f" 성공! 전극 검사기 사양서가 저장되었습니다: {os.path.abspath(output_path)}")
        return output_path
