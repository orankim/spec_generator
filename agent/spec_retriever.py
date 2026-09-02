"""
SpecRetriever — 요구사항(Requirement)을 바탕으로 사내 기존 사양서에서
관련 정보를 검색한다.

기존 generator.py의 `_retrieve_context`는 사용자 문장 전체로 1회
similarity_search만 수행해 슬라이드 단위 chunk를 가져온다. 이 모듈은
같은 Chroma DB(embedding 모델도 동일하게 재사용)를 쓰되:

1. 요구사항의 대상/검사항목/측정원리별로 여러 개의 타겟 질의를 만들어
   각각 검색한 뒤 병합한다 (완전한 슬라이드 대신 항목에 더 가까운 결과를
   모을 수 있다) — retrieve_for_requirement().
2. 선택적으로, PPTX 표의 각 "행"(구분/항목/사양값/비고)을 슬라이드
   전체가 아닌 개별 chunk로도 인덱싱해 진짜 항목 단위 검색이 가능하게
   하는 인덱서를 제공한다 — index_spec_rows_from_folder().
   (기존 build_rag_ollama.py의 슬라이드 단위 인덱싱은 그대로 두고,
   같은 컬렉션에 행 단위 chunk를 "추가"하는 방식이라 기존 기능을
   깨뜨리지 않는다.)
"""
from __future__ import annotations

import os
from glob import glob
from typing import Dict, List, Optional, Tuple

from langchain_community.embeddings import OllamaEmbeddings
from langchain_core.documents import Document
from pptx import Presentation

from . import categorical_match
from .chroma_store import SimpleChromaStore
from .paths import DEFAULT_CHROMA_DB_PATH, resolve_db_path
from .schemas import RequirementSchema
from .units import UnitError, convert, parse_range

_RAG_DEBUG = os.environ.get("RAG_DEBUG", "").lower() in ("1", "true", "yes")

_ITEM_QUERY_HINTS = {
    "thickness": "두께 측정 두께 정확도 두께 분해능",
    "surface_defect": "표면 결함 검출 이물 크랙 핀홀",
    "profile_3d": "3D 형상 프로파일 높이 측정",
    "coating": "코팅 두께 도포량 loading weight",
    "edge_defect": "엣지 결함 가장자리 버 burr",
    # 세부 결함 canonical item(요청서 문제2) — 상위 카테고리(surface_defect/
    # edge_defect)와 별개로 검색 질의를 만들어야 그 세부 항목만 요구된 경우에도
    # 관련 문서가 top-k에서 밀려나지 않는다.
    "scratch": "스크래치 긁힘 표면 결함 검출",
    "contamination": "오염 이물 표면 결함 검출",
    "particle": "파티클 입자 표면 결함 검출",
    "pinhole": "핀홀 표면 결함 검출",
    "void": "보이드 공극 내부 결함",
    "coating_non_uniformity": "코팅 불균일 코팅 두께 편차",
    "edge_crack": "엣지 크랙 가장자리 크랙 균열",
}


def _default_embeddings(host: str) -> OllamaEmbeddings:
    model = os.environ.get("EMBEDDING_MODEL", "bge-m3")
    return OllamaEmbeddings(model=model, base_url=host)


def get_embeddings(host: Optional[str] = None) -> OllamaEmbeddings:
    """
    RAG 구축(build_rag_ollama.py)과 검색(retrieve_for_requirement) 양쪽에서 반드시
    같은 임베딩 모델/서버를 쓰도록 하는 단일 소스. 두 곳이 각자 하드코딩한 값을
    쓰면 벡터 공간이 어긋나 검색이 조용히 실패하므로, 이 함수 하나만 공유한다.
    OLLAMA_HOST/EMBEDDING_MODEL 환경변수를 따른다(.env로 설정 가능).
    """
    resolved_host = host or os.environ.get("OLLAMA_HOST", "http://localhost:11434")
    return _default_embeddings(resolved_host)


def _build_queries(requirement: RequirementSchema) -> List[str]:
    """요구사항에서 여러 개의 타겟 검색 질의를 만든다 (항목 단위 검색)."""
    queries: List[str] = []

    base_terms = []
    if requirement.target.material:
        base_terms.append(requirement.target.material)
    if requirement.measurement_principle:
        base_terms.append(requirement.measurement_principle)
    if requirement.measurement_method:
        base_terms.append(requirement.measurement_method)

    if base_terms:
        queries.append(" ".join(base_terms) + " 검사 설비")

    for item in requirement.inspection_items:
        hint = _ITEM_QUERY_HINTS.get(item, item)
        queries.append(f"{' '.join(base_terms)} {hint}".strip())

    if requirement.raw_text and requirement.raw_text not in queries:
        # material/inspection_items가 이미 있어 의미기반 질의가 만들어지더라도,
        # 사용자 원문(예: "0~200 μm 측정 범위와 ±1 μm 이하 정확도")에 담긴 구체적
        # 수치는 의미기반 질의만으로는 근접하지 않는 문서를 놓칠 수 있다 — 이전에는
        # base_terms/inspection_items가 둘 다 없을 때만 fallback으로 썼는데, 그러면
        # 대부분의 실제 요청에서 raw_text가 전혀 검색에 쓰이지 않는 문제가 있었다.
        queries.append(requirement.raw_text)

    if not queries:
        queries.append("전극 검사 설비 사양")

    return queries


# 세부 결함 이름 등 특정 검사 항목이 실제로 이 chunk에 언급되어 있는지 확인하는
# boost 전용 키워드다. agent.candidate_matcher._INSPECTION_ITEM_DEFECT_KEYWORDS와
# 취지는 같지만(같은 개념), candidate_matcher가 이미 이 모듈(spec_retriever)의
# source_label을 import하므로 반대 방향 import는 순환 import가 된다 — 그래서
# "이 chunk를 검색 결과에 포함시킬지"만 결정하는 최소 사본을 독립적으로 둔다
# (정확한 PASS/FAIL 판정은 여전히 candidate_matcher가 전담).
_ITEM_BOOST_KEYWORDS: Dict[str, Tuple[str, ...]] = {
    "scratch": ("scratch",),
    "contamination": ("contamination", "contaminant"),
    "particle": ("particle",),
    "pinhole": ("pinhole", "pin hole"),
    "void": ("void",),
    "coating_non_uniformity": ("coating non-uniformity", "coating nonuniformity", "coating non uniformity"),
    "edge_crack": ("edge crack",),
}


def _inspection_item_boost_docs(requirement: RequirementSchema, vector_store: SimpleChromaStore) -> List[Document]:
    """
    요구 검사 항목 중 세부 결함 이름(scratch/contamination/particle/pinhole/void/
    coating_non_uniformity/edge_crack)을 실제로 언급하는 chunk를, 의미 검색
    순위와 무관하게 컬렉션 전체에서 찾아온다 — _range_boost_docs와 동일한 원칙
    (구조적으로 확인 가능한 조건은 순위 기반 검색만 믿지 않는다). 이런 세부
    결함 이름은 상위 카테고리("결함")보다 문서 내 등장 빈도가 낮아 의미 검색
    top-k 밖으로 밀려나기 쉽다(실측: "3 μm 이하 크기의 스크래치와 오염" 질의에서
    VisionInspect VI-1000이 top-5 semantic search에서 누락됨).
    """
    keywords: List[str] = []
    for item in requirement.inspection_items:
        keywords.extend(_ITEM_BOOST_KEYWORDS.get(item, ()))
        # profile_3d 등 "Equipment Type/Measurement Principle 서술 텍스트"로 판정하는
        # 항목(agent.categorical_match.INSPECTION_ITEM_CAPABILITY_KEYWORDS)도 같은
        # 이유로 discovery가 밀려날 수 있다 — positive 키워드를 그대로 재사용해
        # boost 대상에 포함시킨다(중복 정의 방지).
        capability_spec = categorical_match.INSPECTION_ITEM_CAPABILITY_KEYWORDS.get(item)
        if capability_spec:
            positive_keywords, _negative_keywords = capability_spec
            keywords.extend(positive_keywords)
    if not keywords:
        return []
    try:
        all_docs = vector_store.get(include=["documents", "metadatas"])
    except Exception:
        return []
    matched: List[Document] = []
    for text, meta in zip(all_docs.get("documents", []) or [], all_docs.get("metadatas", []) or []):
        if not text:
            continue
        text_lower = text.lower()
        if any(kw in text_lower for kw in keywords):
            matched.append(Document(page_content=text, metadata=meta or {}))
    return matched


def _range_boost_docs(requirement: RequirementSchema, vector_store: SimpleChromaStore) -> List[Document]:
    """
    요구사항 원문(raw_text)에 명시된 범위 조건(예: "0~200 μm")을 실제로 포함하는
    문서를, 의미 유사도 순위와 무관하게 컬렉션 전체에서 찾아온다.

    의미기반 벡터 검색은 "비슷한 뜻"은 잘 찾지만 "이 구체적인 수치 범위를 만족하는가"
    같은 구조적 조건은 원래 취약하다 — top-k 밖으로 밀려나면 관련 문서라도 누락된다.
    이 함수는 raw_text에서 파싱한 (min, max, unit)을 각 chunk 원문에서 파싱한
    범위와 직접 비교해(agent.units 재사용, LLM 없음) chunk의 범위가 요구 범위를
    포함하면 결과에 강제로 포함시킨다 — 순위 기반 검색을 대체하는 것이 아니라 보강한다.
    """
    if not requirement.raw_text:
        return []
    range_condition = parse_range(requirement.raw_text)
    if not range_condition:
        return []
    req_lo, req_hi, req_unit = range_condition

    try:
        all_docs = vector_store.get(include=["documents", "metadatas"])
    except Exception:
        return []

    matched: List[Document] = []
    for text, meta in zip(all_docs.get("documents", []) or [], all_docs.get("metadatas", []) or []):
        if not text:
            continue
        doc_range = parse_range(text)
        if not doc_range:
            continue
        d_lo, d_hi, d_unit = doc_range
        try:
            d_lo_c = convert(d_lo, d_unit, req_unit)
            d_hi_c = convert(d_hi, d_unit, req_unit)
        except UnitError:
            continue
        if d_lo_c <= req_lo and d_hi_c >= req_hi:
            matched.append(Document(page_content=text, metadata=meta or {}))
    return matched


class RetrievedChunk(Document):
    """langchain Document와 동일 — 타입 힌트 가독성을 위한 별칭."""


def retrieve_for_requirement(
    requirement: RequirementSchema,
    db_path: Optional[str] = None,
    ollama_host: Optional[str] = None,
    k_per_query: int = 15,
) -> List[Document]:
    """
    요구사항 기반 다중 질의 검색을 수행하고, (source, content) 기준으로
    중복 제거한 Document 목록을 반환한다.

    k_per_query 기본값 15(이전 10) — agent/pipeline.py retrieve_and_generate()의
    docstring에 근거(실제 bge-m3 임베딩 기반 k=[5,10,15,20] 전체 56케이스 재현
    실험) 기록. k=10에서는 Retrieval Recall이 86.0%(6/43 MISS)였고, MISS 6건
    전부 순위 경쟁(정답 문서가 실제로는 순위 11~19위로 검색됐으나 top-10 밖으로
    밀림)으로 확인됐다 — k=15에서 Recall 97.7%(42/43, MISS 6건 중 5건 해소),
    No-Match 안전성(False PASS 0건)은 그대로 유지됨을 실측으로 확인한 뒤 올렸다.

    db_path를 명시하지 않으면 CHROMA_DB_PATH 환경변수 -> 저장소 루트 기준 기본값
    (agent/paths.DEFAULT_CHROMA_DB_PATH) 순으로 정해진다 — build_rag_ollama.py도
    동일한 규칙을 쓰므로, 두 프로세스를 서로 다른 작업 디렉터리에서 실행해도
    (예: 빌드는 프로젝트 루트에서, 서버는 다른 cwd에서) 항상 같은 디스크 경로를
    가리킨다. 예전에는 둘 다 "./chroma_db_specs"라는 *상대경로* 기본값을 썼는데,
    이는 각자의 cwd 기준으로 따로 해석되어 서로 다른(한쪽은 비어있는) 디렉터리를
    가리킬 수 있었다 — 예외 없이 조용히 검색 결과 0개로 이어지는 원인이었다.
    """
    resolved_db_path = resolve_db_path(db_path)
    host = ollama_host or os.environ.get("OLLAMA_HOST", "http://localhost:11434")
    embeddings = _default_embeddings(host)
    vector_store = SimpleChromaStore(persist_directory=resolved_db_path, embedding_function=embeddings)

    queries = _build_queries(requirement)
    seen = set()
    results: List[Document] = []
    for query in queries:
        hits_with_score = vector_store.similarity_search_with_score(query, k=k_per_query)
        for doc, dist in hits_with_score:
            score = 1.0 / (1.0 + dist)
            doc.metadata["score"] = score
            key = (doc.metadata.get("source"), doc.page_content[:80])
            if key in seen:
                continue
            seen.add(key)
            results.append(doc)
        if _RAG_DEBUG:
            print(f"[RAG DEBUG] query: {query!r} -> {len(hits_with_score)}개 hit (db_path={resolved_db_path})")

    for doc in _range_boost_docs(requirement, vector_store):
        key = (doc.metadata.get("source"), doc.page_content[:80])
        if key in seen:
            continue
        seen.add(key)
        results.append(doc)
        if _RAG_DEBUG:
            print(f"[RAG DEBUG] range_boost -> 추가됨: {source_label(doc)!r}")

    for doc in _inspection_item_boost_docs(requirement, vector_store):
        key = (doc.metadata.get("source"), doc.page_content[:80])
        if key in seen:
            continue
        seen.add(key)
        results.append(doc)
        if _RAG_DEBUG:
            print(f"[RAG DEBUG] inspection_item_boost -> 추가됨: {source_label(doc)!r}")

    # 후보 확정 후 전체 문서 로드 (실사용자 보고 버그): 초기 top-k 의미 검색이
    # 문서 자체는 찾아내면서도(=이 문서가 후보라는 것은 확정됨) 그 문서의 특정
    # chunk(예: "## Inspection Target"의 Maximum Electrode Width, "Measurement
    # Performance"의 Measurement Speed)는 우연히 top-k 밖으로 밀려나는 경우가
    # 실제로 있었다 — candidate_matcher._extract_candidate_fact는 여기 넘겨준
    # chunk만 볼 수 있으므로, 문서에 값이 있는데도 UNKNOWN으로 잘못 판정되는
    # 원인이 됐다(예: SPEC-013 Width, SPEC-039 Speed). RAG 검색(top-k)의 역할은
    # "어떤 문서가 후보인가"를 정하는 데까지만이고, 일단 후보로 확정되면 그
    # 문서 전체(모든 chunk)를 deterministic field extraction의 입력으로 삼는다
    # (요청서 구조: RAG 검색 → 후보 선정 → 후보 전체 문서 로드 → deterministic
    # extraction → Hard Requirement 검증). _pull_identity_chunks(chunk_id=0만)의
    # 상위 호환이므로 그 함수는 더 이상 쓰지 않는다.
    matched_filenames = {doc.metadata.get("filename") for doc in results if doc.metadata.get("filename")}
    for doc in _pull_full_candidate_documents(vector_store, matched_filenames):
        key = (doc.metadata.get("source"), doc.page_content[:80])
        if key in seen:
            continue
        seen.add(key)
        results.append(doc)
        if _RAG_DEBUG:
            print(f"[RAG DEBUG] full_candidate_doc -> 추가됨: {source_label(doc)!r}")

    if _RAG_DEBUG:
        try:
            collection_count = vector_store._collection.count()
        except Exception:
            collection_count = "?"
        print(
            f"[RAG DEBUG] collection: {vector_store._collection.name} | "
            f"document_count: {collection_count} | queries: {len(queries)} | "
            f"retrieved_chunks(중복제거 후): {len(results)}"
        )
        for i, doc in enumerate(results, start=1):
            print(f"[RAG DEBUG]   [{i}] source: {source_label(doc)}")
            print(f"[RAG DEBUG]       content: {doc.page_content[:120]!r}")

    return results


def _pull_full_candidate_documents(vector_store: SimpleChromaStore, filenames: set) -> List[Document]:
    """
    확정된 후보 문서(파일)마다 컬렉션에 있는 모든 chunk를 조건 없이 통째로
    가져온다(의미 검색 순위와 무관) — Manufacturer/Model이 적힌 "## General"
    chunk부터 Width/Speed가 적힌 "## Inspection Target"/"## Measurement
    Performance" chunk까지 전부 포함되므로, 이전에 chunk_id=0(식별 정보)만
    가져오던 것의 상위 호환이다. PPTX 소스(filename 메타데이터 없음)에 대해서는
    filenames 집합 자체에 들어오지 않아(retrieve_for_requirement에서 filter됨)
    조용히 건너뛴다.
    """
    full_docs: List[Document] = []
    for filename in filenames:
        if not filename:
            continue
        try:
            raw = vector_store.get(where={"filename": filename}, include=["documents", "metadatas"])
        except Exception:
            continue
        for text, meta in zip(raw.get("documents", []) or [], raw.get("metadatas", []) or []):
            full_docs.append(Document(page_content=text, metadata=meta or {}))
    return full_docs


def source_label(doc: Document) -> str:
    """
    참고 자료 출처를 사람이 읽을 짧은 형태로 만든다. 항상 짧은 파일명(filename)을
    우선하고, 없으면 source(과거 PPTX 문서는 filename 없이 source만 있을 수 있음)로
    폴백한다 — "sample_specs/spec_01.md" 같은 전체 경로가 아니라 "spec_01.md"만
    사용자/LLM에게 노출한다 (요청서 11절).
    """
    return doc.metadata.get("filename") or doc.metadata.get("source", "Unknown")


def format_context(docs: List[Document]) -> str:
    """검색 결과를 프롬프트에 넣을 수 있는 텍스트로 변환한다. Markdown(category/item)과
    PPTX(slide_number) 두 출처 형식을 모두 지원한다."""
    parts = []
    for i, doc in enumerate(docs, start=1):
        source = source_label(doc)
        if doc.metadata.get("source_type") == "markdown":
            location = doc.metadata.get("category") or doc.metadata.get("item") or ""
            if doc.metadata.get("item"):
                location = f"{doc.metadata.get('category', '')} > {doc.metadata['item']}"
            header = f"{source} ({location})" if location else source
        else:
            slide = doc.metadata.get("slide_number", "?")
            header = f"{source} Slide {slide}"
        parts.append(f"\n[참고 자료 {i} (출처: {header})]\n{doc.page_content}\n")
    return "".join(parts)


# ==========================================
# 항목(행) 단위 인덱싱 — 선택적 보강
# ==========================================
def parse_pptx_rows(file_path: str) -> List[Document]:
    """
    PPTX 표의 각 행(구분/항목/사양값/비고)을 개별 Document로 만든다.
    슬라이드 전체를 chunk로 삼는 build_rag_ollama.parse_pptx_file과 달리,
    "정확도가 얼마인가" 같은 항목 단위 질의에 더 가깝게 매칭되는 chunk를 만든다.
    """
    prs = Presentation(file_path)
    file_name = os.path.basename(file_path)
    documents: List[Document] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            headers = [cell.text.strip() for cell in table.rows[0].cells]
            for row in list(table.rows)[1:]:
                cells = [cell.text.strip() for cell in row.cells]
                if not any(cells):
                    continue
                row_dict = dict(zip(headers, cells))
                category = row_dict.get("구분", "")
                item = row_dict.get("항목", "")
                value = row_dict.get("사양값", "")
                note = row_dict.get("비고", "")
                content = f"[{category}] {item}: {value}" + (f" ({note})" if note else "")
                documents.append(
                    Document(
                        page_content=content,
                        metadata={
                            "source": file_name,
                            "slide_number": slide_idx,
                            "category": category,
                            "item": item,
                            "chunk_type": "spec_row",
                        },
                    )
                )
    return documents


def index_spec_rows_from_folder(
    pptx_folder: str,
    db_path: Optional[str] = None,
    ollama_host: Optional[str] = None,
) -> int:
    """
    폴더 내 모든 PPTX의 표 행을 파싱해 기존 Chroma 컬렉션에 "추가"한다.
    (기존 슬라이드 단위 chunk는 그대로 두고 보강하는 방식.)
    반환값: 추가된 chunk 수.
    """
    host = ollama_host or os.environ.get("OLLAMA_HOST", "http://localhost:11434")
    embeddings = _default_embeddings(host)

    pptx_files = glob(os.path.join(pptx_folder, "*.pptx"))
    all_rows: List[Document] = []
    for f in pptx_files:
        all_rows.extend(parse_pptx_rows(f))

    if not all_rows:
        return 0

    vector_store = SimpleChromaStore(persist_directory=resolve_db_path(db_path), embedding_function=embeddings)
    vector_store.add_documents(all_rows)
    return len(all_rows)
