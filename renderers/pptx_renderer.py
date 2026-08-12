"""
Specification JSON -> PPTX 렌더러.

PPTX는 이제 "여러 출력 포맷 중 하나"로 취급한다. 템플릿 파일(회사 지정 양식)이
있으면 그것을 사용하고, 없으면 템플릿 없이도 python-pptx로 기본 PPTX를 바로
생성한다 — 회사 템플릿을 git에 올릴 수 없는 사내망 환경에서도 Agent 파이프라인
전체(스키마 검증/Markdown/HTML 생성 포함)가 항상 동작해야 하기 때문이다.

- 템플릿 있음: 기존 agent/pptx_electrode_builder.py(ElectrodeSpecPPTXBuilder)를
  그대로 재사용한다 (수정하지 않음 — 기존 PPTX Generator 유지 원칙).
- 템플릿 없음: renderers/common.py의 동일한 RenderSection 모델로 기본 PPTX를
  즉석에서 만든다 (Markdown/HTML과 같은 데이터 소스를 사용하므로 세 포맷 간
  내용이 어긋나지 않는다).
"""
from __future__ import annotations

import os
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from agent.schemas import ComplianceRecord, RequirementSchema, SpecificationSchema, ValidationResult
from agent.spec_validator import build_compliance_report

from .common import (
    COMPLIANCE_SECTION_EMPTY_NOTE,
    COMPLIANCE_SECTION_TITLE,
    RenderSection,
    build_notes_section,
    build_sections,
    build_validation_section,
)

_NUMERIC_SECTION_IDS = {
    "inspection_target",
    "inspection_requirements",
    "measurement_performance",
    "spatial_performance",
    "inspection_performance",
    "defect_inspection",
}


def _default_template_path() -> Optional[str]:
    """
    회사 템플릿 경로. 코드에 하드코딩하지 않고 환경변수로만 받는다 — 실제 회사
    템플릿 파일은 git에 커밋하지 않는 것이 원칙이므로, 저장소 안의 고정 경로를
    기본값으로 두지 않는다.
    """
    return os.environ.get("PPT_TEMPLATE_PATH") or None


def _add_compliance_slide(prs: Presentation, blank, records):
    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = COMPLIANCE_SECTION_TITLE
    tp.font.size = Pt(22)
    tp.font.bold = True

    if not records:
        note_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(9), Inches(1))
        note_box.text_frame.text = COMPLIANCE_SECTION_EMPTY_NOTE
        return slide

    headers = ["Item", "Unit", "Requirement", "Specification", "Result", "Reason"]
    table_shape = slide.shapes.add_table(len(records) + 1, len(headers), Inches(0.4), Inches(1.2), Inches(9.2), Inches(5.8))
    table = table_shape.table
    for c, header in enumerate(headers):
        table.rows[0].cells[c].text = header
    for r_idx, record in enumerate(records, start=1):
        values = [
            record.item,
            record.unit or "-",
            "UNKNOWN" if record.requirement is None else str(record.requirement),
            "UNKNOWN" if record.specification is None else str(record.specification),
            record.result,
            record.reason,
        ]
        for c, val in enumerate(values):
            cell = table.rows[r_idx].cells[c]
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
    return slide


def _add_data_table_slide(prs: Presentation, blank, section: RenderSection, headers, row_to_values):
    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = section.title
    tp.font.size = Pt(22)
    tp.font.bold = True

    if not section.rows:
        note_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(9), Inches(1))
        note_box.text_frame.text = section.note or "No data."
        return slide

    n_rows = len(section.rows) + 1
    table_shape = slide.shapes.add_table(n_rows, len(headers), Inches(0.4), Inches(1.2), Inches(9.2), Inches(5.8))
    table = table_shape.table
    for c, header in enumerate(headers):
        table.rows[0].cells[c].text = header

    for r, row in enumerate(section.rows, start=1):
        values = row_to_values(row)
        for c, val in enumerate(values):
            cell = table.rows[r].cells[c]
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
    return slide


def _build_default_presentation(
    specification: SpecificationSchema,
    requirement: Optional[RequirementSchema],
    validation: Optional[ValidationResult],
    title: str,
) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    title_box = cover.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    sub = cover.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1.5))
    sub.text_frame.text = specification.equipment.name or "(장비명 미정)"
    sub.text_frame.paragraphs[0].font.size = Pt(18)

    for section in build_sections(specification):
        if section.id in _NUMERIC_SECTION_IDS:
            headers = ["Item", "Unit", "Specification", "Status", "Source"]
            row_to_values = lambda row: [row.label, row.unit or "-", row.value_display, row.status or "-", row.source or "-"]
        else:
            headers = ["Item", "Specification"]
            row_to_values = lambda row: [row.label, row.value_display]
        _add_data_table_slide(prs, blank, section, headers, row_to_values)

    _add_data_table_slide(
        prs, blank, build_validation_section(validation),
        ["Level / Field", "Message"], lambda row: [row.label, row.value_display],
    )

    compliance_records = build_compliance_report(specification, requirement)
    _add_compliance_slide(prs, blank, compliance_records)

    _add_data_table_slide(
        prs, blank, build_notes_section(specification),
        ["Item", "Specification"], lambda row: [row.label, row.value_display],
    )

    return prs


def render_pptx(
    specification: SpecificationSchema,
    output_path: str,
    requirement: Optional[RequirementSchema] = None,
    validation: Optional[ValidationResult] = None,
    template_path: Optional[str] = None,
    title: str = "Electrode Inspection Equipment Specification",
) -> str:
    """
    Specification -> PPTX. template_path가 주어지지 않으면 PPT_TEMPLATE_PATH
    환경변수를 확인하고, 그마저 없으면(또는 파일이 실제로 없으면) 템플릿 없이
    기본 PPTX를 생성한다 — 어떤 경우에도 예외 없이 PPTX 한 장은 나온다.
    """
    resolved_template = template_path or _default_template_path()

    if resolved_template and os.path.exists(resolved_template):
        # 기존 PPTX Generator 재사용 (수정하지 않음)
        from agent.pptx_electrode_builder import ElectrodeSpecPPTXBuilder

        builder = ElectrodeSpecPPTXBuilder(template_path=resolved_template)
        return builder.build(specification, output_path=output_path)

    prs = _build_default_presentation(specification, requirement, validation, title)
    prs.save(output_path)
    return output_path
